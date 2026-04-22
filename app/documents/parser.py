"""
Document parsing utilities for PDF, PPTX, and TXT files.
"""

import fitz  # PyMuPDF
from pptx import Presentation
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


def parse_pdf(file_path: Path) -> str:
    """Extract text from a PDF file using PyMuPDF."""
    try:
        doc = fitz.open(str(file_path))
        text_parts = []
        for page_num, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{text}")
        doc.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PDF parsing error: {e}")
        raise ValueError(f"Failed to parse PDF: {e}")


def parse_pptx(file_path: Path) -> str:
    """Extract text from a PowerPoint file."""
    try:
        prs = Presentation(str(file_path))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PPTX parsing error: {e}")
        raise ValueError(f"Failed to parse PPTX: {e}")


def parse_txt(file_path: Path) -> str:
    """Read text from a plain text file."""
    try:
        return file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return file_path.read_text(encoding="latin-1")


def parse_document(file_path: Path, file_type: str) -> str:
    """Parse a document based on its file type. Returns extracted text."""
    parsers = {
        "pdf": parse_pdf,
        "pptx": parse_pptx,
        "txt": parse_txt,
    }

    parser = parsers.get(file_type.lower())
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")

    text = parser(file_path)

    if not text or not text.strip():
        raise ValueError("No text content could be extracted from the document")

    return text
